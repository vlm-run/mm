import re

from docx import Document
from pptx import Presentation
from pptx.shapes.base import BaseShape


def extract_docx(path: str) -> str:
    """Extract text from a DOCX file, preserving some structure."""
    doc = Document(path)
    chunks = []

    for block in doc.element.body:
        tag = block.tag.split("}")[-1]  # strip namespace

        if tag == "p":
            # paragraph — includes headings
            from docx.text.paragraph import Paragraph

            para = Paragraph(block, doc)
            if para.text.strip():
                style_name = para.style.name if para.style and para.style.name else ""
                prefix = f"[{style_name}] " if "Heading" in style_name else ""
                chunks.append(f"{prefix}{para.text}\n")

        elif tag == "tbl":
            from docx.table import Table

            table = Table(block, doc)
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                chunks.append(" | ".join(cells))

    return "\n".join(chunks)


def extract_pptx(path: str) -> str:
    """Extract text from a PPTX file"""
    prs = Presentation(path)
    chunks = []

    for slide_num, slide in enumerate(prs.slides, 1):
        chunks.append(f"--- [Slide {slide_num}. LAYOUT: {slide.slide_layout.name}] ---")

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            if slide.notes_slide.notes_text_frame.text.strip():
                chunks.append(f"[Notes] {slide.notes_slide.notes_text_frame.text.strip()}")

        text_chunks = []
        image_chunks = []

        def get_shape_content(shape: BaseShape):
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                sorted_subshapes = sorted(
                    shape.shapes,
                    key=lambda x: (
                        float("-inf") if not x.top else x.top,
                        float("-inf") if not x.left else x.left,
                    ),
                )
                for subshape in sorted_subshapes:
                    get_shape_content(subshape)

            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if text := para.text.strip():
                        text_chunks.append(text)

            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    text_chunks.append(" | ".join(cells))

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "") or shape.name
                except Exception:
                    alt_text = shape.name
                alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text).strip()
                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                image_chunks.append(f"![{alt_text}]({filename})")

        sorted_shapes = sorted(
            slide.shapes,
            key=lambda x: (
                float("-inf") if not x.top else x.top,
                float("-inf") if not x.left else x.left,
            ),
        )
        for shape in sorted_shapes:
            get_shape_content(shape)

        chunks.extend(text_chunks)
        chunks.extend(image_chunks)

    return "\n".join(chunks)
